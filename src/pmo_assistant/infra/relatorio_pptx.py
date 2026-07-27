"""Geração do Status Report em PPTX a partir de um Cronograma em memória.

Sempre em memória (io.BytesIO) — Streamlit Cloud tem filesystem efêmero, nunca
escreva em disco aqui. Fonte dos dados é sempre st.session_state.cronograma
(ver core/relatorio.py) — nunca listar_tarefas_por_projeto(), que descarta as
linhas-resumo persistidas no banco.

# TODO: Curva S (percentual acumulado por mês) — fora de escopo do MVP.
"""

from __future__ import annotations

import io
from datetime import date
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from pmo_assistant.core.relatorio import (
    atividades_realizadas,
    classificar_tarefas_por_status,
    pontos_de_atencao,
    proximas_atividades,
    tarefa_raiz,
)

if TYPE_CHECKING:
    from collections.abc import Callable

    from pptx.presentation import Presentation as PresentationType
    from pptx.slide import Slide

    from pmo_assistant.core.models import Cronograma, TarefaCronograma

# ---------------------------------------------------------------------------
# Paleta — identidade visual própria do relatório, independente do tema verde
# do PMO Assistente em si (o PPTX é um deliverable para o cliente final).
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x25, 0x45)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
CINZA_CLARO = RGBColor(0xF4, 0xF6, 0xF8)
CINZA_LINHA = RGBColor(0xE8, 0xEC, 0xEF)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TEXTO_ESCURO = RGBColor(0x1B, 0x1F, 0x23)
VERMELHO_ATRASO = RGBColor(0xE6, 0x39, 0x46)
AMARELO_ALERTA = RGBColor(0xF2, 0xC9, 0x4C)
FONTE = "Calibri"

# ---------------------------------------------------------------------------
# Geometria (slide 16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGEM_X = Inches(0.4)
ALTURA_HEADER = Inches(0.9)
ALTURA_UNDERLINE = Inches(0.08)

LARGURA_CARD = Inches(2.34)
GAP_CARD = Inches(0.2)
TOPO_CARDS = Inches(1.7)
ALTURA_CARD = Inches(1.3)

TOPO_GRAFICO = Inches(3.3)
LARGURA_GRAFICO = SLIDE_WIDTH - 2 * MARGEM_X
ALTURA_GRAFICO = Inches(3.3)

TOPO_TABELA = Inches(1.3)
ALTURA_LINHA = Inches(0.4)
LARGURA_TABELA = Inches(12.3)
MAX_LINHAS_TABELA = 12

TOPO_ESTADO_VAZIO = Inches(3.2)

MAX_CHARS_NOME = 80


# ---------------------------------------------------------------------------
# Helpers de formatação
# ---------------------------------------------------------------------------
def _fmt_data(d: date | None) -> str:
    return d.strftime("%d/%m/%Y") if d else "N/D"


def _fmt_pct(p: float | None) -> str:
    return f"{p:.0f}%" if p is not None else "N/D"


def _truncar(texto: str, limite: int = MAX_CHARS_NOME) -> str:
    return texto if len(texto) <= limite else texto[: limite - 1].rstrip() + "…"


# ---------------------------------------------------------------------------
# Helpers de shape
# ---------------------------------------------------------------------------
def _sem_borda_sem_sombra(shape) -> None:
    shape.line.fill.background()
    shape.shadow.inherit = False


def _estilizar_run(run, *, tamanho: int, cor: RGBColor, negrito: bool = False) -> None:
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    run.font.color.rgb = cor
    run.font.name = FONTE


def _adicionar_cabecalho(slide: Slide, titulo: str) -> None:
    """Barra NAVY + underline TEAL + título — usado nos slides 2 a 5."""
    barra = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, ALTURA_HEADER
    )
    _sem_borda_sem_sombra(barra)
    barra.fill.solid()
    barra.fill.fore_color.rgb = NAVY

    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), ALTURA_HEADER, SLIDE_WIDTH, ALTURA_UNDERLINE
    )
    _sem_borda_sem_sombra(underline)
    underline.fill.solid()
    underline.fill.fore_color.rgb = TEAL

    caixa = slide.shapes.add_textbox(
        MARGEM_X, Inches(0.15), SLIDE_WIDTH - 2 * MARGEM_X, Inches(0.6)
    )
    tf = caixa.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = titulo
    _estilizar_run(run, tamanho=28, cor=BRANCO, negrito=True)


def _texto_vazio(slide: Slide, mensagem: str, cor: RGBColor = TEXTO_ESCURO) -> None:
    """Estado vazio centralizado — reusado pelas 3 tabelas sem dados e pelo
    gráfico do slide 2 quando todas as categorias estão zeradas (mesma faixa
    vertical, mesma necessidade visual)."""
    caixa = slide.shapes.add_textbox(
        MARGEM_X, TOPO_ESTADO_VAZIO, SLIDE_WIDTH - 2 * MARGEM_X, Inches(1)
    )
    tf = caixa.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = mensagem
    _estilizar_run(run, tamanho=18, cor=cor)


def _tabela_atividades(
    slide: Slide,
    tarefas: list[TarefaCronograma],
    colunas: tuple[str, str, str],
    valor_col2: Callable[[TarefaCronograma], str],
    valor_col3: Callable[[TarefaCronograma], str],
    cor_texto: RGBColor,
) -> None:
    """Tabela com header NAVY/branco e linhas alternando BRANCO/CINZA_LINHA."""
    linhas = len(tarefas) + 1
    tabela_shape = slide.shapes.add_table(
        linhas, 3, MARGEM_X, TOPO_TABELA, LARGURA_TABELA, ALTURA_LINHA * linhas
    )
    tabela = tabela_shape.table
    tabela.columns[0].width = Inches(7.3)
    tabela.columns[1].width = Inches(2.5)
    tabela.columns[2].width = Inches(2.5)

    for col, titulo_coluna in enumerate(colunas):
        celula = tabela.cell(0, col)
        celula.text = titulo_coluna
        celula.fill.solid()
        celula.fill.fore_color.rgb = NAVY
        run = celula.text_frame.paragraphs[0].runs[0]
        _estilizar_run(run, tamanho=12, cor=BRANCO, negrito=True)

    for i, t in enumerate(tarefas, start=1):
        fundo = BRANCO if i % 2 == 1 else CINZA_LINHA
        valores = (_truncar(t.nome), valor_col2(t), valor_col3(t))
        for col, valor in enumerate(valores):
            celula = tabela.cell(i, col)
            celula.text = valor
            celula.fill.solid()
            celula.fill.fore_color.rgb = fundo
            run = celula.text_frame.paragraphs[0].runs[0]
            _estilizar_run(run, tamanho=11, cor=cor_texto)


def _nota_excedente(slide: Slide, total: int, exibidos: int) -> None:
    excedente = total - exibidos
    if excedente <= 0:
        return
    topo = TOPO_TABELA + ALTURA_LINHA * (exibidos + 1) + Inches(0.05)
    caixa = slide.shapes.add_textbox(MARGEM_X, topo, LARGURA_TABELA, Inches(0.3))
    p = caixa.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f"+{excedente} atividades adicionais não exibidas"
    _estilizar_run(run, tamanho=10, cor=TEXTO_ESCURO)


def _slide_com_tabela(
    prs: PresentationType,
    titulo: str,
    tarefas: list[TarefaCronograma],
    colunas: tuple[str, str, str],
    valor_col2: Callable[[TarefaCronograma], str],
    valor_col3: Callable[[TarefaCronograma], str],
    mensagem_vazia: str,
    *,
    cor_vazia: RGBColor = TEXTO_ESCURO,
    cor_texto_linhas: RGBColor = TEXTO_ESCURO,
    ordenar_por: Callable[[TarefaCronograma], object] | None = None,
) -> None:
    """Monta um slide 'lista de atividades' completo — usado pelos slides 3,
    4 e 5, que só diferem em título, filtro de tarefas, colunas e cor."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _adicionar_cabecalho(slide, titulo)

    if not tarefas:
        _texto_vazio(slide, mensagem_vazia, cor=cor_vazia)
        return

    if ordenar_por is not None:
        tarefas = sorted(tarefas, key=ordenar_por, reverse=True)
    exibidos = tarefas[:MAX_LINHAS_TABELA]
    _tabela_atividades(slide, exibidos, colunas, valor_col2, valor_col3, cor_texto_linhas)
    _nota_excedente(slide, len(tarefas), len(exibidos))


# ---------------------------------------------------------------------------
# Slide 1 — Capa
# ---------------------------------------------------------------------------
def _slide_capa(
    prs: PresentationType, cronograma: Cronograma, cliente: str | None, gerado_em: date
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY

    caixa_titulo = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.6), SLIDE_WIDTH - Inches(1), Inches(1.2)
    )
    tf = caixa_titulo.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = cronograma.nome_projeto
    _estilizar_run(run, tamanho=36, cor=BRANCO, negrito=True)

    topo_status = Inches(3.9)
    if cliente is not None:
        caixa_cliente = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.9), SLIDE_WIDTH - Inches(1), Inches(0.5)
        )
        p_cliente = caixa_cliente.text_frame.paragraphs[0]
        p_cliente.alignment = PP_ALIGN.CENTER
        run_cliente = p_cliente.add_run()
        run_cliente.text = f"Cliente: {cliente}"
        _estilizar_run(run_cliente, tamanho=18, cor=TEAL)
        topo_status = Inches(4.4)

    caixa_status = slide.shapes.add_textbox(
        Inches(0.5), topo_status, SLIDE_WIDTH - Inches(1), Inches(0.5)
    )
    p_status = caixa_status.text_frame.paragraphs[0]
    p_status.alignment = PP_ALIGN.CENTER
    run_status = p_status.add_run()
    run_status.text = f"Status Report — {gerado_em.strftime('%d/%m/%Y')}"
    _estilizar_run(run_status, tamanho=20, cor=BRANCO)

    caixa_rodape = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), SLIDE_WIDTH - Inches(1), Inches(0.4)
    )
    p_rodape = caixa_rodape.text_frame.paragraphs[0]
    p_rodape.alignment = PP_ALIGN.CENTER
    run_rodape = p_rodape.add_run()
    run_rodape.text = "Gerado por PMO Assistente"
    _estilizar_run(run_rodape, tamanho=10, cor=CINZA_CLARO)


# ---------------------------------------------------------------------------
# Slide 2 — Resumo Geral
# ---------------------------------------------------------------------------
def _card(slide: Slide, x, titulo: str, valor: str) -> None:
    caixa = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, TOPO_CARDS, LARGURA_CARD, ALTURA_CARD)
    _sem_borda_sem_sombra(caixa)
    caixa.fill.solid()
    caixa.fill.fore_color.rgb = CINZA_CLARO

    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)

    p_label = tf.paragraphs[0]
    p_label.alignment = PP_ALIGN.CENTER
    run_label = p_label.add_run()
    run_label.text = titulo.upper()
    _estilizar_run(run_label, tamanho=10, cor=NAVY, negrito=True)

    p_valor = tf.add_paragraph()
    p_valor.alignment = PP_ALIGN.CENTER
    run_valor = p_valor.add_run()
    run_valor.text = valor
    _estilizar_run(run_valor, tamanho=22, cor=TEXTO_ESCURO, negrito=True)


def _grafico_status(slide: Slide, contagem: dict[str, int]) -> None:
    categorias = ["Concluída", "No Prazo", "Atrasada", "Tarefa Futura"]
    cores = [TEAL, NAVY, VERMELHO_ATRASO, CINZA_LINHA]

    dados = CategoryChartData()
    dados.categories = categorias
    dados.add_series("Tarefas", tuple(contagem[c] for c in categorias))

    quadro = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        MARGEM_X,
        TOPO_GRAFICO,
        LARGURA_GRAFICO,
        ALTURA_GRAFICO,
        dados,
    )
    grafico = quadro.chart
    grafico.has_legend = False
    plotagem = grafico.plots[0]
    plotagem.vary_by_categories = True
    plotagem.has_data_labels = True
    plotagem.data_labels.font.size = Pt(12)
    plotagem.data_labels.font.name = FONTE
    serie = plotagem.series[0]
    for ponto, cor in zip(serie.points, cores, strict=False):
        ponto.format.fill.solid()
        ponto.format.fill.fore_color.rgb = cor
    grafico.category_axis.tick_labels.font.size = Pt(11)
    grafico.value_axis.tick_labels.font.size = Pt(11)


def _slide_resumo_geral(prs: PresentationType, cronograma: Cronograma) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _adicionar_cabecalho(slide, "Resumo Geral")

    raiz = tarefa_raiz(cronograma)
    cartoes = [
        ("% Concluído Previsto", _fmt_pct(raiz.percentual_esperado if raiz else None)),
        ("% Concluído Real", _fmt_pct(raiz.percentual_concluido if raiz else None)),
        ("Início Linha de Base", _fmt_data(raiz.inicio_baseline if raiz else None)),
        ("Término Linha de Base", _fmt_data(raiz.termino_baseline if raiz else None)),
        ("Término (Impacto)", _fmt_data(raiz.termino if raiz else None)),
    ]
    for i, (titulo_card, valor_card) in enumerate(cartoes):
        x = MARGEM_X + i * (LARGURA_CARD + GAP_CARD)
        _card(slide, x, titulo_card, valor_card)

    contagem = classificar_tarefas_por_status(cronograma.tarefas)
    if sum(contagem.values()) == 0:
        _texto_vazio(slide, "Sem dados de tarefas para exibir.")
    else:
        _grafico_status(slide, contagem)


# ---------------------------------------------------------------------------
# Slides 3, 4, 5 — listas de atividades
# ---------------------------------------------------------------------------
def _slide_atividades_realizadas(prs: PresentationType, cronograma: Cronograma) -> None:
    _slide_com_tabela(
        prs,
        "Atividades Realizadas",
        atividades_realizadas(cronograma.tarefas),
        ("Atividade", "Entrega Prevista", "Entrega Real"),
        lambda t: _fmt_data(t.termino_baseline),
        lambda t: _fmt_data(t.termino_real),
        "Nenhuma atividade concluída neste período.",
        ordenar_por=lambda t: t.termino_real or date.min,
    )


def _slide_proximas_atividades(prs: PresentationType, cronograma: Cronograma) -> None:
    _slide_com_tabela(
        prs,
        "Próximas Atividades",
        proximas_atividades(cronograma.tarefas),
        ("Atividade", "Entrega Prevista", "Nova Previsão"),
        lambda t: _fmt_data(t.termino_baseline),
        lambda t: _fmt_data(t.termino),
        "Nenhuma atividade pendente registrada.",
    )


def _slide_pontos_de_atencao(prs: PresentationType, cronograma: Cronograma) -> None:
    _slide_com_tabela(
        prs,
        "Pontos de Atenção",
        pontos_de_atencao(cronograma.tarefas),
        ("Atividade", "Prevista", "Atual"),
        lambda t: _fmt_data(t.termino_baseline),
        lambda t: _fmt_data(t.termino),
        "Nenhuma tarefa atrasada — projeto dentro do previsto.",
        cor_vazia=TEAL,
        cor_texto_linhas=VERMELHO_ATRASO,
    )


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------
def gerar_status_report_pptx(
    cronograma: Cronograma, cliente: str | None, gerado_em: date
) -> bytes:
    """Gera o Status Report em PPTX (5 slides) a partir do Cronograma em memória.

    Retorna bytes prontos para st.download_button — não escreve em disco
    (Streamlit Cloud tem filesystem efêmero). O cronograma deve vir de
    st.session_state (ver core/relatorio.py) — reconstruído do banco perde
    as linhas-resumo, essenciais para o card "Resumo Geral" (slide 2).
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _slide_capa(prs, cronograma, cliente, gerado_em)
    _slide_resumo_geral(prs, cronograma)
    _slide_atividades_realizadas(prs, cronograma)
    _slide_proximas_atividades(prs, cronograma)
    _slide_pontos_de_atencao(prs, cronograma)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
