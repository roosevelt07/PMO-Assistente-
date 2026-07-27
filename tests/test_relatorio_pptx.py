"""Smoke tests de infra/relatorio_pptx.py — confirma que o PPTX é gerado e
abre sem exceção nos casos de borda (sem linha-resumo, listas vazias,
categorias zeradas)."""

from __future__ import annotations

import io
from datetime import date

from pptx import Presentation

from pmo_assistant.core.models import Cronograma, TarefaCronograma
from pmo_assistant.infra.relatorio_pptx import gerar_status_report_pptx


def _abrir(pptx_bytes: bytes) -> Presentation:
    return Presentation(io.BytesIO(pptx_bytes))


def test_gera_cinco_slides_com_dados_completos():
    cr = Cronograma(
        projeto_id=1,
        nome_projeto="Charqueadas TLC",
        data_referencia=date(2026, 7, 1),
        tarefas=[
            TarefaCronograma(
                id_tarefa=1,
                nome="Projeto",
                percentual_concluido=94.0,
                percentual_esperado=99.0,
                inicio_baseline=date(2025, 12, 3),
                termino_baseline=date(2026, 7, 9),
                termino=date(2026, 8, 11),
                eh_resumo=True,
            ),
            TarefaCronograma(
                id_tarefa=2,
                nome="Concluída",
                percentual_concluido=100.0,
                termino_baseline=date(2026, 5, 1),
                termino_real=date(2026, 5, 3),
            ),
            TarefaCronograma(
                id_tarefa=3,
                nome="Atrasada",
                percentual_concluido=40.0,
                percentual_esperado=80.0,
                termino_baseline=date(2026, 6, 1),
                termino=date(2026, 7, 1),
            ),
            TarefaCronograma(
                id_tarefa=4,
                nome="Futura",
                percentual_concluido=0.0,
                termino_baseline=date(2026, 9, 1),
            ),
        ],
    )
    pptx_bytes = gerar_status_report_pptx(cr, "Grantel-Axia", date(2026, 7, 26))
    prs = _abrir(pptx_bytes)
    assert len(prs.slides) == 5


def test_sem_tarefa_raiz_mostra_nd_sem_excecao():
    cr = Cronograma(
        projeto_id=1,
        nome_projeto="Sem Resumo",
        tarefas=[TarefaCronograma(id_tarefa=2, nome="Sub", percentual_concluido=50.0)],
    )
    pptx_bytes = gerar_status_report_pptx(cr, None, date(2026, 7, 26))
    prs = _abrir(pptx_bytes)
    assert len(prs.slides) == 5


def test_listas_vazias_nao_lanca_excecao():
    cr = Cronograma(projeto_id=1, nome_projeto="Vazio", tarefas=[])
    pptx_bytes = gerar_status_report_pptx(cr, None, date(2026, 7, 26))
    prs = _abrir(pptx_bytes)
    assert len(prs.slides) == 5


def test_categorias_zeradas_pula_grafico_sem_excecao():
    # única tarefa é eh_resumo=True -> classificar_tarefas_por_status soma 0
    cr = Cronograma(
        projeto_id=1,
        nome_projeto="Só Resumo",
        tarefas=[
            TarefaCronograma(
                id_tarefa=1, nome="Resumo", percentual_concluido=50.0, eh_resumo=True
            )
        ],
    )
    pptx_bytes = gerar_status_report_pptx(cr, None, date(2026, 7, 26))
    prs = _abrir(pptx_bytes)
    assert len(prs.slides) == 5


def test_cliente_none_gera_menos_shapes_na_capa():
    """Confirma que o bloco condicional 'Cliente: {cliente}' realmente é
    omitido quando cliente is None (checagem estrutural)."""
    cr = Cronograma(projeto_id=1, nome_projeto="P", tarefas=[])
    sem_cliente = _abrir(gerar_status_report_pptx(cr, None, date(2026, 7, 26)))
    com_cliente = _abrir(gerar_status_report_pptx(cr, "ACME", date(2026, 7, 26)))
    assert len(sem_cliente.slides[0].shapes) < len(com_cliente.slides[0].shapes)
